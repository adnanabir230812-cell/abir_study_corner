import os
import sys
from pptx import Presentation
from pypdf import PdfReader

# Paths
SRC_DIR = r"E:\Abir study\3.2 1st Part\3.2 1st Part\AT-3201 Principles of Plant Pathology and Diseases of Field Crops\Rezaul Sir"
OUT_DIR = r"C:\Users\Lenovo\.gemini\antigravity\scratch\toppr\extracted_raw_secb"
UPLOADS_DIR = r"C:\Users\Lenovo\.gemini\antigravity\scratch\toppr\uploads"

# Create folders
os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(UPLOADS_DIR, exist_ok=True)

def sanitize_filename(name):
    return "".join(c for c in name if c.isalnum() or c in (' ', '_', '-', '.')).strip()

def extract_pptx(filepath, out_txt_path):
    print(f"Extracting PPTX: {filepath}")
    prs = Presentation(filepath)
    base_name = os.path.splitext(os.path.basename(filepath))[0]
    sanitized_base = sanitize_filename(base_name)
    
    with open(out_txt_path, "w", encoding="utf-8") as f:
        f.write(f"# Course Slides: {base_name}\n\n")
        
        for slide_idx, slide in enumerate(prs.slides):
            f.write(f"## Slide {slide_idx + 1}\n\n")
            
            # Sort shapes by top and left coordinates to read text top-to-bottom, left-to-right
            shapes = []
            for shape in slide.shapes:
                top = getattr(shape, "top", 0)
                if top is None:
                    top = 0
                left = getattr(shape, "left", 0)
                if left is None:
                    left = 0
                shapes.append((top, left, shape))
            
            # Since top and left could be Emu objects, we can convert them to standard ints for sorting
            shapes.sort(key=lambda x: (int(x[0]), int(x[1])))
            
            img_idx = 1
            for _, _, shape in shapes:
                # 1. Text
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    f.write(f"{text}\n\n")
                
                # 2. Tables
                if shape.has_table:
                    table = shape.table
                    f.write("\n")
                    for row_idx, row in enumerate(table.rows):
                        row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        f.write("| " + " | ".join(row_data) + " |\n")
                        if row_idx == 0:
                            f.write("|" + "---|" * len(row_data) + "\n")
                    f.write("\n")
                
                # 3. Pictures
                if shape.shape_type == 13 or (shape.is_placeholder and hasattr(shape, "image")):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        ext = image.ext or "png"
                        
                        img_filename = f"{sanitized_base}_s{slide_idx + 1}_img{img_idx}.{ext}"
                        img_path = os.path.join(UPLOADS_DIR, img_filename)
                        
                        with open(img_path, "wb") as img_f:
                            img_f.write(image_bytes)
                        
                        f.write(f"\n![Diagram: Slide {slide_idx + 1} Image {img_idx}](/uploads/{img_filename})\n\n")
                        img_idx += 1
                    except Exception as e:
                        pass
            
            f.write("\n---\n\n")

def extract_pdf(filepath, out_txt_path):
    print(f"Extracting PDF: {filepath}")
    reader = PdfReader(filepath)
    base_name = os.path.splitext(os.path.basename(filepath))[0]
    sanitized_base = sanitize_filename(base_name)
    
    with open(out_txt_path, "w", encoding="utf-8") as f:
        f.write(f"# Course Document: {base_name}\n\n")
        
        for page_idx, page in enumerate(reader.pages):
            f.write(f"## Page {page_idx + 1}\n\n")
            text = page.extract_text()
            if text:
                f.write(text.strip())
                f.write("\n\n")
            
            # Extract images from PDF page if any
            img_idx = 1
            if "/Resources" in page and "/XObject" in page["/Resources"]:
                xobjects = page["/Resources"]["/XObject"].get_object()
                for obj_name in xobjects:
                    xobject = xobjects[obj_name]
                    if xobject["/Subtype"] == "/Image":
                        try:
                            data = xobject.get_data()
                            ext = "png"
                            if "/Filter" in xobject:
                                if xobject["/Filter"] == "/DCTDecode":
                                    ext = "jpg"
                                elif xobject["/Filter"] == "/JPXDecode":
                                    ext = "jp2"
                            
                            img_filename = f"{sanitized_base}_p{page_idx + 1}_img{img_idx}.{ext}"
                            img_path = os.path.join(UPLOADS_DIR, img_filename)
                            
                            with open(img_path, "wb") as img_f:
                                img_f.write(data)
                            
                            f.write(f"\n![Diagram: Page {page_idx + 1} Image {img_idx}](/uploads/{img_filename})\n\n")
                            img_idx += 1
                        except Exception as e:
                            pass
            
            f.write("\n---\n\n")

def main():
    files_in_dir = os.listdir(SRC_DIR)
    print(f"Found {len(files_in_dir)} files in source directory.")
    
    for fn in files_in_dir:
        ext = os.path.splitext(fn)[1].lower()
        if ext not in ('.pptx', '.pdf'):
            print(f"Skipping non-supported file: {fn}")
            continue
            
        full_path = os.path.join(SRC_DIR, fn)
        out_txt_filename = sanitize_filename(os.path.splitext(fn)[0]) + ".md"
        out_txt_path = os.path.join(OUT_DIR, out_txt_filename)
        
        try:
            if ext == ".pptx":
                extract_pptx(full_path, out_txt_path)
            elif ext == ".pdf":
                extract_pdf(full_path, out_txt_path)
        except Exception as e:
            print(f"Error processing {fn}: {e}")

if __name__ == "__main__":
    main()

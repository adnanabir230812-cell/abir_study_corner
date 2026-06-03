import os
import sys
from pptx import Presentation
from pypdf import PdfReader

# Paths
SRC_DIR = r"E:\Abir study\3.2 1st Part\3.2 1st Part\AT-3201 Principles of Plant Pathology and Diseases of Field Crops\Sabiha Madame\21 Batch"
OUT_DIR = r"C:\Users\Lenovo\.gemini\antigravity\scratch\toppr\extracted_raw"
UPLOADS_DIR = r"C:\Users\Lenovo\.gemini\antigravity\scratch\toppr\uploads"

# Create folders
os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(UPLOADS_DIR, exist_ok=True)

# Selected files
FILES_TO_PROCESS = [
    "1.0 Introduction and History of Plant Pathology.pptx",
    "1.1 Pathogenesis.pptx",
    "1.2 Toxin.pptx",
    "1.3 Effects of Pathogens on Plant Physiological Functions.pptx",
    "2.1 Dessimination of the Pathogen.pptx",
    "2.2 2.2 Epidemiology(IMPORTANT).pdf",  # Wait, let's verify exact name. It was "2.2 Epidemiology(IMPORTANT).pdf" in directory list
    "2.2 Epidemiology(IMPORTANT).pdf",
    "2.2 Epidemiology.pptx",
    "3.1 Method of Plant Disease Control.pptx",
    "3.2 Biological Control.pptx",
    "3.3 Chemical Control.pptx"
]

# Ensure we use unique names and don't duplicate
FILES_TO_PROCESS = list(set(FILES_TO_PROCESS))

def sanitize_filename(name):
    return "".join(c for c in name if c.isalnum() or c in (' ', '_', '-', '.')).strip()

def extract_pptx(filepath, out_txt_path):
    print(f"Extracting PPTX: {filepath}")
    prs = Presentation(filepath)
    base_name = os.path.splitext(os.path.basename(filepath))[0]
    sanitized_base = sanitize_filename(base_name)
    
    with open(out_txt_path, "w", encoding="utf-8") as f:
        f.write(f"# Course Slides: {base_name}\n\n")
        
        for slide_idx, slide in enumerate(prs.slides):
            f.write(f"## Slide {slide_idx + 1}\n\n")
            
            # Sort shapes by top and left coordinates to read text top-to-bottom, left-to-right
            shapes = []
            for shape in slide.shapes:
                # Add shape and its positions
                top = getattr(shape, "top", 0)
                left = getattr(shape, "left", 0)
                shapes.append((top, left, shape))
            
            shapes.sort(key=lambda x: (x[0], x[1]))
            
            img_idx = 1
            for _, _, shape in shapes:
                # 1. Text
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # If it looks like a list, keep it
                    f.write(f"{text}\n\n")
                
                # 2. Tables
                if shape.has_table:
                    table = shape.table
                    f.write("\n")
                    for row_idx, row in enumerate(table.rows):
                        row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        f.write("| " + " | ".join(row_data) + " |\n")
                        if row_idx == 0:
                            f.write("|" + "---|" * len(row_data) + "\n")
                    f.write("\n")
                
                # 3. Pictures
                # Shape type 13 is Picture, shape type 17 is Placeholder (which might contain picture)
                if shape.shape_type == 13 or (shape.is_placeholder and hasattr(shape, "image")):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        ext = image.ext or "png"
                        
                        img_filename = f"{sanitized_base}_s{slide_idx + 1}_img{img_idx}.{ext}"
                        img_path = os.path.join(UPLOADS_DIR, img_filename)
                        
                        with open(img_path, "wb") as img_f:
                            img_f.write(image_bytes)
                        
                        # Add image reference to markdown
                        # Web route serves this at /uploads/filename
                        f.write(f"\n![Diagram: Slide {slide_idx + 1} Image {img_idx}](/uploads/{img_filename})\n\n")
                        img_idx += 1
                    except Exception as e:
                        print(f"Error saving image on slide {slide_idx + 1}: {e}")
            
            f.write("\n---\n\n")

def extract_pdf(filepath, out_txt_path):
    print(f"Extracting PDF: {filepath}")
    reader = PdfReader(filepath)
    base_name = os.path.splitext(os.path.basename(filepath))[0]
    sanitized_base = sanitize_filename(base_name)
    
    with open(out_txt_path, "w", encoding="utf-8") as f:
        f.write(f"# Course Document: {base_name}\n\n")
        
        for page_idx, page in enumerate(reader.pages):
            f.write(f"## Page {page_idx + 1}\n\n")
            text = page.extract_text()
            if text:
                f.write(text.strip())
                f.write("\n\n")
            
            # Extract images from PDF page if any
            img_idx = 1
            if "/Resources" in page and "/XObject" in page["/Resources"]:
                xobjects = page["/Resources"]["/XObject"].get_object()
                for obj_name in xobjects:
                    xobject = xobjects[obj_name]
                    if xobject["/Subtype"] == "/Image":
                        try:
                            # Try to extract the image
                            data = xobject.get_data()
                            # Check extension
                            ext = "png"
                            if "/Filter" in xobject:
                                if xobject["/Filter"] == "/DCTDecode":
                                    ext = "jpg"
                                elif xobject["/Filter"] == "/JPXDecode":
                                    ext = "jp2"
                            
                            img_filename = f"{sanitized_base}_p{page_idx + 1}_img{img_idx}.{ext}"
                            img_path = os.path.join(UPLOADS_DIR, img_filename)
                            
                            with open(img_path, "wb") as img_f:
                                img_f.write(data)
                            
                            f.write(f"\n![Diagram: Page {page_idx + 1} Image {img_idx}](/uploads/{img_filename})\n\n")
                            img_idx += 1
                        except Exception as e:
                            print(f"Error saving PDF image on page {page_idx + 1}: {e}")
            
            f.write("\n---\n\n")

def main():
    files_in_dir = os.listdir(SRC_DIR)
    print("Files found in source directory:")
    for fn in files_in_dir:
        print(f"- {fn}")
        
    for target_fn in FILES_TO_PROCESS:
        # Find exact case-insensitive match in files_in_dir
        matched_fn = None
        for fn in files_in_dir:
            if fn.lower() == target_fn.lower():
                matched_fn = fn
                break
                
        if not matched_fn:
            print(f"WARNING: Target file '{target_fn}' not found in directory. Skipping.")
            continue
            
        full_path = os.path.join(SRC_DIR, matched_fn)
        out_txt_filename = sanitize_filename(os.path.splitext(matched_fn)[0]) + ".md"
        out_txt_path = os.path.join(OUT_DIR, out_txt_filename)
        
        try:
            if matched_fn.lower().endswith(".pptx"):
                extract_pptx(full_path, out_txt_path)
            elif matched_fn.lower().endswith(".pdf"):
                extract_pdf(full_path, out_txt_path)
            else:
                print(f"Unknown extension: {matched_fn}")
        except Exception as e:
            print(f"Error processing {matched_fn}: {e}")

if __name__ == "__main__":
    main()
